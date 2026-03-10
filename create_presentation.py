"""
Gera apresentacao PowerPoint v2 — versao melhorada com:
- Fontes/referencias em slides com dados
- Termos tecnicos + linguagem acessivel
- Requisitos minimos VR, claustrofobia, escala
- Ciencia da imersao
- Gaussian Splatting: historia e ciencia
- Definicao de Game Engine
- Apps de referencia para criacao VR
- Glossario tecnico A-Frame
- Explicacao HTML mais detalhada
- QR codes para referencia
- Exemplos de ferramentas AR/VR faceis de usar
"""

import os, io, tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import qrcode

# ── Cores do tema ──────────────────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x11, 0x17)
BG_CARD    = RGBColor(0x16, 0x1B, 0x22)
ACCENT     = RGBColor(0x58, 0xA6, 0xFF)   # azul
ACCENT2    = RGBColor(0x3F, 0xB9, 0x50)   # verde
ACCENT3    = RGBColor(0xF7, 0x81, 0x66)   # coral
ACCENT4    = RGBColor(0xD2, 0xA8, 0xFF)   # roxo
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC4)
MID_GRAY   = RGBColor(0x6E, 0x76, 0x81)
DARK_GRAY  = RGBColor(0x44, 0x4C, 0x56)
SOURCE_CLR = RGBColor(0x55, 0x5F, 0x6B)   # cinza para fontes
TBL_HDR    = RGBColor(0x1A, 0x2A, 0x3A)
TBL_ROW1   = RGBColor(0x12, 0x18, 0x20)
TBL_ROW2   = RGBColor(0x16, 0x1E, 0x28)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── QR code helper ─────────────────────────────────────────────
qr_temp_files = []

def make_qr(url, size=4):
    qr = qrcode.QRCode(version=1, box_size=10, border=2)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="white", back_color="#0D1117")
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    img.save(tmp.name)
    qr_temp_files.append(tmp.name)
    return tmp.name

# ── Helpers ────────────────────────────────────────────────────
def dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=18, color=WHITE, bold=False,
             alignment=PP_ALIGN.LEFT, space_before=Pt(6),
             font_name="Segoe UI", italic=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p

def add_bullet_list(slide, left, top, width, height, items,
                    font_size=18, color=WHITE, bullet_char="\u2022 ",
                    line_spacing=Pt(8)):
    tf = add_textbox(slide, left, top, width, height,
                     bullet_char + items[0], font_size, color)
    for item in items[1:]:
        add_para(tf, bullet_char + item, font_size, color,
                 space_before=line_spacing)
    return tf

def add_card(slide, left, top, width, height, fill_color=BG_CARD,
             border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_source(slide, text, y=Inches(7.0)):
    """Adiciona referencia/fonte no fundo do slide."""
    tf = add_textbox(slide, Inches(0.8), y, Inches(11.5), Inches(0.4),
                     text, 10, SOURCE_CLR)
    tf.paragraphs[0].font.italic = True

def title_slide(title, subtitle="", notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(1), Inches(2.6), Inches(2), Pt(4))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
    add_textbox(slide, Inches(1), Inches(2.9), Inches(11), Inches(1.5),
                title, 40, WHITE, True, font_name="Segoe UI Semibold")
    if subtitle:
        add_textbox(slide, Inches(1), Inches(4.3), Inches(11), Inches(1),
                    subtitle, 22, LIGHT_GRAY, font_name="Segoe UI Light")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def section_slide(section_num, section_title, subtitle="", notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    add_textbox(slide, Inches(1), Inches(1.5), Inches(2), Inches(2),
                str(section_num).zfill(2), 96, ACCENT, True, font_name="Segoe UI Light")
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(1), Inches(3.7), Inches(1.5), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
    add_textbox(slide, Inches(1), Inches(4), Inches(11), Inches(1.2),
                section_title, 36, WHITE, True, font_name="Segoe UI Semibold")
    if subtitle:
        add_textbox(slide, Inches(1), Inches(5.2), Inches(10), Inches(1),
                    subtitle, 20, LIGHT_GRAY)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def content_slide(title, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                title, 28, WHITE, True, font_name="Segoe UI Semibold")
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.15), Inches(1.2), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_table(slide, left, top, width, height, rows, cols,
              headers, data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = "Segoe UI"
        cell.fill.solid()
        cell.fill.fore_color.rgb = TBL_HDR
    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = LIGHT_GRAY
                paragraph.font.name = "Segoe UI"
            cell.fill.solid()
            cell.fill.fore_color.rgb = TBL_ROW1 if i % 2 == 0 else TBL_ROW2
    return tbl


# ══════════════════════════════════════════════════════════════
#  SLIDE 1 — CAPA
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), H)
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
add_textbox(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(0.5),
            "MESTRADO EM SAUDE  |  AULA PRATICA  |  3 HORAS", 14, ACCENT)
add_textbox(slide, Inches(1.2), Inches(2.3), Inches(10), Inches(1.8),
            "Realidade Estendida e\nMetaverso em Saude", 48, WHITE, True,
            font_name="Segoe UI Semibold")
add_textbox(slide, Inches(1.2), Inches(4.2), Inches(10), Inches(0.8),
            "Ferramentas de Criacao de Conteudo", 28, ACCENT, font_name="Segoe UI Light")
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(5.4), Inches(3), Pt(3))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
add_textbox(slide, Inches(1.2), Inches(5.7), Inches(10), Inches(0.5),
            "2025/2026", 16, MID_GRAY, font_name="Segoe UI Light")
slide.notes_slide.notes_text_frame.text = "Slide de abertura."


# ══════════════════════════════════════════════════════════════
#  SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════
slide = content_slide("Agenda da Aula")
blocks = [
    ("01", "Fundamentos de XR", "Imersao, cybersickness, requisitos\ntecnicos, exemplos clinicos", "40 min", ACCENT),
    ("02", "Ferramentas de Conteudo", "Som, imagem 360, video, 3D,\nGaussian Splatting, Kiri Engine", "40 min", ACCENT2),
    ("03", "Interatividade", "Game Engines, Unity, Unreal,\nA-Frame, Inspector", "25 min", ACCENT3),
    ("04", "Exercicio Pratico", "Codigo A-Frame guiado,\ntestar em VR no Quest", "50 min", ACCENT4),
]
for idx, (num, title, desc, dur, color) in enumerate(blocks):
    x = Inches(0.8 + idx * 3.05)
    add_card(slide, x, Inches(1.8), Inches(2.8), Inches(4.2), BG_CARD, color)
    add_textbox(slide, x + Inches(0.3), Inches(2.1), Inches(2.2), Inches(0.8),
                num, 42, color, True, font_name="Segoe UI Light")
    add_textbox(slide, x + Inches(0.3), Inches(2.9), Inches(2.2), Inches(0.6),
                title, 18, WHITE, True)
    add_textbox(slide, x + Inches(0.3), Inches(3.6), Inches(2.2), Inches(1.2),
                desc, 14, LIGHT_GRAY)
    add_textbox(slide, x + Inches(0.3), Inches(5.3), Inches(2.2), Inches(0.4),
                dur, 14, color, True)


# ══════════════════════════════════════════════════════════════
#  BLOCO 1 — FUNDAMENTOS XR
# ══════════════════════════════════════════════════════════════
section_slide("01", "Fundamentos de XR e Imersao",
    "Realidade Estendida, ciencia da imersao, requisitos tecnicos")

# ── O que e XR ──
slide = content_slide("O que e Realidade Estendida (XR)?",
    "Explicar o espectro com exemplos do dia-a-dia.")
items = [
    "AR (Augmented Reality / Realidade Aumentada) — sobrepoe informacao digital ao mundo real",
    "MR (Mixed Reality / Realidade Mista) — objetos digitais interagem com o mundo real",
    "VR (Virtual Reality / Realidade Virtual) — substitui totalmente o ambiente por um digital",
    "XR (Extended Reality / Realidade Estendida) — termo guarda-chuva: AR + MR + VR"
]
add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(7), Inches(3.5), items, 17, LIGHT_GRAY)

add_card(slide, Inches(8.5), Inches(1.6), Inches(4.2), Inches(5), BG_CARD, ACCENT)
spectrum = [
    ("Mundo Real", "100% fisico", MID_GRAY),
    ("AR", "Digital SOBRE real", ACCENT2),
    ("MR", "Digital + real INTERAGEM", ACCENT3),
    ("VR", "100% digital (imersao total)", ACCENT),
]
for i, (label, desc, color) in enumerate(spectrum):
    yy = Inches(1.9 + i * 1.1)
    add_textbox(slide, Inches(9), yy, Inches(3.5), Inches(0.4), label, 20, color, True)
    add_textbox(slide, Inches(9), yy + Inches(0.35), Inches(3.5), Inches(0.4), desc, 14, LIGHT_GRAY)

tf = add_textbox(slide, Inches(0.8), Inches(5.5), Inches(7), Inches(1.5),
    "Definicao tecnica (IEEE):", 14, ACCENT, True)
add_para(tf, '"XR refers to all real-and-virtual combined environments and human-machine', 12, SOURCE_CLR, italic=True)
add_para(tf, 'interactions generated by computer technology and wearables." — IEEE Std 2048-2020', 12, SOURCE_CLR, italic=True)

# ── Ferramentas AR/VR faceis ──
slide = content_slide("Ferramentas AR e VR Faceis de Usar",
    "Mostrar que existem ferramentas acessiveis sem programacao.")

add_card(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.2), BG_CARD, ACCENT)
tf = add_textbox(slide, Inches(1.1), Inches(1.8), Inches(5.3), Inches(0.5),
    "AR — Realidade Aumentada", 20, ACCENT, True)
items_ar = [
    "Adobe Aero — criar AR sem codigo (drag & drop)",
    "Zappar — AR web-based, criacao visual",
    "Assemblr Studio — AR educativo, interface simples",
    "Reality Composer (Apple) — AR para iOS nativo",
    "Instagram / TikTok Filters — AR social (Spark AR, Lens Studio)",
    "Google ARCore Geospatial — AR com geolocalizacao",
    "Merge EDU — cubos AR para educacao medica",
]
for item in items_ar:
    add_para(tf, "\u2022 " + item, 14, LIGHT_GRAY, space_before=Pt(5))

add_card(slide, Inches(7.1), Inches(1.6), Inches(5.8), Inches(5.2), BG_CARD, ACCENT2)
tf = add_textbox(slide, Inches(7.4), Inches(1.8), Inches(5.3), Inches(0.5),
    "VR — Realidade Virtual", 20, ACCENT2, True)
items_vr = [
    "A-Frame — VR com HTML no browser (usamos hoje!)",
    "CoSpaces Edu — VR educativo drag & drop",
    "Mozilla Hubs — salas VR sociais, sem instalar",
    "Spatial.io — espacos VR colaborativos",
    "Gravity Sketch — modelacao 3D DENTRO do VR",
    "Tilt Brush / Open Brush — pintar em 3D no VR",
    "Shapes XR — prototipagem VR dentro do headset",
    "Meta Horizon Worlds — criar mundos VR sem codigo",
]
for item in items_vr:
    add_para(tf, "\u2022 " + item, 14, LIGHT_GRAY, space_before=Pt(5))


# ── Porque importa na saude ──
slide = content_slide("Porque Importa na Saude?")
stats = [("$2.1B", "2019"), ("$30.4B", "2026 (previsao)"), ("42.4%", "CAGR")]
for i, (val, label) in enumerate(stats):
    x = Inches(0.8 + i * 4)
    add_card(slide, x, Inches(1.6), Inches(3.5), Inches(2), BG_CARD, ACCENT)
    add_textbox(slide, x + Inches(0.3), Inches(1.8), Inches(3), Inches(1), val, 44, ACCENT, True)
    add_textbox(slide, x + Inches(0.3), Inches(2.8), Inches(3), Inches(0.5), label, 16, LIGHT_GRAY)

tf = add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(2.5),
    "Mercado global de VR em saude (Healthcare VR Market)", 20, WHITE, True)
add_para(tf, "\u2022 Treino cirurgico com simulacao haptica e feedback em tempo real", 16, LIGHT_GRAY)
add_para(tf, "\u2022 Educacao anatomica imersiva (modelos 3D interativos)", 16, LIGHT_GRAY)
add_para(tf, "\u2022 Terapia da dor: RelieVRx — 1.o dispositivo VR aprovado pela FDA (2021)", 16, LIGHT_GRAY)
add_para(tf, "\u2022 Reabilitacao gamificada (neuroplasticidade + motivacao)", 16, LIGHT_GRAY)
add_source(slide, "Fontes: Grand View Research, 2023; FDA De Novo Classification DEN200033 (RelieVRx); ITIF, 2025")


# ── Exemplos reais ──
slide = content_slide("Exemplos Clinicos: Treino Cirurgico")
tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
    "Evidencia Cientifica", 22, ACCENT, True)
add_para(tf, "", 6, WHITE)
add_para(tf, "Osso VR — Treino Ortopedico (Orthopedic Surgical Simulation)", 18, WHITE, True)
add_para(tf, "  Estudo com 38 estudantes: grupo VR superou grupo tradicional", 16, LIGHT_GRAY)
add_para(tf, "  na proficiencia processual (p < 0.001), percentagem de passos", 16, LIGHT_GRAY)
add_para(tf, "  corretos (p < 0.002) e conhecimento de instrumentos (p < 0.01)", 16, LIGHT_GRAY)
add_para(tf, "", 6, WHITE)
add_para(tf, "FundamentalVR — Simulacao Haptica", 18, WHITE, True)
add_para(tf, "  Melhoria de 44% na precisao cirurgica vs metodos tradicionais", 16, LIGHT_GRAY)
add_para(tf, "  Metricas: forca aplicada, precisao processual, eficiencia temporal", 16, LIGHT_GRAY)
add_para(tf, "", 6, WHITE)
add_para(tf, "UNC-Chapel Hill (2025) — Seguranca no Bloco Operatorio", 18, WHITE, True)
add_para(tf, "  30 minutos de treino VR melhoraram 90% dos comportamentos de seguranca", 16, LIGHT_GRAY)
add_para(tf, "  observados (study published in peer-reviewed journal)", 16, LIGHT_GRAY)
add_source(slide, "Fontes: Osso VR Clinical Studies (ossovr.com); FundamentalVR white papers; UNC-Chapel Hill, Jan 2026 (unc.edu/posts/2026/01/27)")


slide = content_slide("Exemplos Clinicos: Dor, Fobias, Reabilitacao")
add_card(slide, Inches(0.8), Inches(1.6), Inches(3.6), Inches(4.5), BG_CARD, ACCENT3)
tf = add_textbox(slide, Inches(1.1), Inches(1.8), Inches(3), Inches(0.5),
                 "Gestao da Dor", 20, ACCENT3, True)
add_para(tf, "RelieVRx (AppliedVR)", 14, WHITE, True, space_before=Pt(10))
add_para(tf, "FDA De Novo: DEN200033", 12, SOURCE_CLR, italic=True)
add_para(tf, "66% reducao >30% dor", 14, LIGHT_GRAY, space_before=Pt(8))
add_para(tf, "46% atingiram >50%", 14, LIGHT_GRAY)
add_para(tf, "Efeito mantido 3 meses", 14, LIGHT_GRAY)

add_card(slide, Inches(4.8), Inches(1.6), Inches(3.6), Inches(4.5), BG_CARD, ACCENT4)
tf = add_textbox(slide, Inches(5.1), Inches(1.8), Inches(3), Inches(0.5),
                 "Fobias (VRET)", 20, ACCENT4, True)
add_para(tf, "VR Exposure Therapy", 14, WHITE, True, space_before=Pt(10))
add_para(tf, "Taxa sucesso: 66-90%", 14, LIGHT_GRAY, space_before=Pt(8))
add_para(tf, "3% recusa VR vs 27% in-vivo", 14, LIGHT_GRAY)
add_para(tf, "Dropout: 16% < CBT", 14, LIGHT_GRAY)
add_para(tf, "Alturas: -68% medo medio", 14, LIGHT_GRAY)

add_card(slide, Inches(8.8), Inches(1.6), Inches(3.6), Inches(4.5), BG_CARD, ACCENT2)
tf = add_textbox(slide, Inches(9.1), Inches(1.8), Inches(3), Inches(0.5),
                 "Reabilitacao", 20, ACCENT2, True)
add_para(tf, "Cochrane Review 2025", 14, WHITE, True, space_before=Pt(10))
add_para(tf, "Membro sup.: SMD=0.42", 14, LIGHT_GRAY, space_before=Pt(8))
add_para(tf, "Equilibrio: SMD=0.68", 14, LIGHT_GRAY)
add_para(tf, "So 31% cumprem fisioterapia", 14, LIGHT_GRAY)
add_para(tf, "VR gamificada resolve adesao", 14, LIGHT_GRAY)

add_source(slide, "Fontes: Frontiers in VR 2025 (RelieVRx); Frontiers in Psychology 2019 (VRET meta-analysis); Cochrane Library CD008349.pub5 (2025)")


slide = content_slide("Educacao Anatomica em VR")
stats_anat = [
    ("75%", "dos estudos comparativos:\nVR superou metodos\ntradicionais (cadaveres,\naulas, modelos 2D)"),
    ("71%", "dos estudos mostram\nmelhoria estatisticamente\nsignificativa na\naprendizagem"),
    ("100%", "dos estudos de\npercecao: satisfacao\nunanimemente\nfavoravel"),
]
for i, (val, label) in enumerate(stats_anat):
    x = Inches(0.8 + i * 4)
    add_card(slide, x, Inches(1.6), Inches(3.5), Inches(3.5), BG_CARD, ACCENT)
    add_textbox(slide, x + Inches(0.3), Inches(1.9), Inches(3), Inches(1), val, 52, ACCENT, True, PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.3), Inches(3.1), Inches(3), Inches(1.5), label, 14, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.8),
    "Particularmente eficaz: neuroanatomia, sistema musculoesqueletico, anatomia cardiaca", 16, MID_GRAY)
add_source(slide, "Fontes: Adnan et al., BMC Medical Education (2025); Systematic Review PMC12051093 (2025); JMIR Medical Education e62803")


# ── Ciencia da Imersao ──
slide = content_slide("A Ciencia da Imersao: Porque Funciona?",
    "Explicar com base neurocientifica. Os medicos apreciam a base cientifica.")

tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(5.5),
    "Tres pilares neurocientificos", 20, ACCENT, True)
add_para(tf, "", 6, WHITE)
add_para(tf, "1. PRESENCA (Presence)", 18, WHITE, True)
add_para(tf, "   Sensacao psicologica de 'estar la'. O cerebro aceita o ambiente", 15, LIGHT_GRAY)
add_para(tf, "   virtual como real. Mediada pelo cortex parietal posterior e", 15, LIGHT_GRAY)
add_para(tf, "   insula anterior. Mensuravel por questionarios (IPQ, SUS).", 15, LIGHT_GRAY)
add_para(tf, "", 6, WHITE)
add_para(tf, "2. IMERSAO (Immersion)", 18, WHITE, True)
add_para(tf, "   Capacidade tecnica do sistema de estimular os sentidos:", 15, LIGHT_GRAY)
add_para(tf, "   \u2022 Visual: estereoscopia, FOV (field of view), resolucao", 15, LIGHT_GRAY)
add_para(tf, "   \u2022 Auditivo: audio espacial (HRTF), ambisonics", 15, LIGHT_GRAY)
add_para(tf, "   \u2022 Haptico: feedback tatil (controladores, luvas)", 15, LIGHT_GRAY)
add_para(tf, "   \u2022 Propriocetivo: hand tracking, body tracking", 15, LIGHT_GRAY)
add_para(tf, "", 6, WHITE)
add_para(tf, "3. AGENCIA (Agency)", 18, WHITE, True)
add_para(tf, "   Capacidade de agir e ver consequencias no mundo virtual.", 15, LIGHT_GRAY)
add_para(tf, "   Ativa o loop sensorio-motor: acao \u2192 feedback \u2192 aprendizagem.", 15, LIGHT_GRAY)

add_card(slide, Inches(7.5), Inches(1.6), Inches(5.2), Inches(5.2), BG_CARD, ACCENT)
tf = add_textbox(slide, Inches(7.8), Inches(1.8), Inches(4.6), Inches(0.5),
    "O que faz funcionar?", 18, ACCENT, True)
items = [
    "Embodiment: sensacao de ter um corpo virtual (rubber hand illusion)",
    "Congruencia multisensorial: visao + audio + tato alinhados",
    "Affordances: objetos virtuais 'convidam' a acao",
    "Emotional engagement: reacoes emocionais reais a estimulos virtuais",
    "Flow state: imersao total na tarefa (Csikszentmihalyi)",
    "Memoria espacial: retencao superior com contexto 3D vs 2D",
]
for item in items:
    add_para(tf, "\u2022 " + item, 13, LIGHT_GRAY, space_before=Pt(6))

add_source(slide, "Fontes: Slater & Sanchez-Vives, 2016 (Enhancing Our Lives with Immersive VR); Riva et al., 2007 (Affective Interactions Using VR)")


# ── Cybersickness ──
slide = content_slide("Cybersickness: O Conflito Sensorial")
tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(5),
    "Etiologia: Sensory Conflict Theory (Reason & Brand, 1975)", 18, ACCENT, True)
add_para(tf, "", 6, WHITE)
add_para(tf, "Conflito Vestibulo-Ocular (inter-sensorial):", 16, WHITE, True)
add_para(tf, "Os olhos reportam movimento (fluxo optico), mas o sistema", 15, LIGHT_GRAY)
add_para(tf, "vestibular (canais semicirculares + otolitos) reporta ausencia", 15, LIGHT_GRAY)
add_para(tf, "de movimento. O cerebro interpreta como intoxicacao \u2192 nausea.", 15, LIGHT_GRAY)
add_para(tf, "", 8, WHITE)
add_para(tf, "Sintomas (SSQ — Simulator Sickness Questionnaire):", 16, WHITE, True)
add_para(tf, "\u2022 Nausea, tonturas, desorientacao", 15, LIGHT_GRAY)
add_para(tf, "\u2022 Fadiga ocular (oculomotor strain)", 15, LIGHT_GRAY)
add_para(tf, "\u2022 Sudorese, cefaleia", 15, LIGHT_GRAY)
add_para(tf, "\u2022 Prevalencia: 20-80% dependendo do conteudo", 15, LIGHT_GRAY)
add_para(tf, "", 8, WHITE)
add_para(tf, "Claustrofobia e VR:", 16, ACCENT3, True)
add_para(tf, "Headsets cobrem o campo visual \u2192 pode desencadear", 15, LIGHT_GRAY)
add_para(tf, "claustrofobia. Solucao: passthrough mode, sessoes curtas,", 15, LIGHT_GRAY)
add_para(tf, "ambientes abertos, permitir remover headset a qualquer momento.", 15, LIGHT_GRAY)

add_card(slide, Inches(7.5), Inches(1.6), Inches(5.2), Inches(5), BG_CARD, ACCENT3)
tf = add_textbox(slide, Inches(7.8), Inches(1.8), Inches(4.6), Inches(0.5),
    "Estrategias de Mitigacao", 18, ACCENT3, True)
items = [
    "90+ fps (ver requisitos tecnicos)",
    "Motion-to-photon < 20ms",
    "Vinheta (FOV restriction) em locomocao",
    "Teletransporte vs locomocao continua",
    "Pontos referencia fixos (cockpit, nariz)",
    "Sessoes: iniciar 5-10 min, aumentar",
    "Posicao sentada reduz conflito",
    "NUNCA tirar controlo da camara",
    "Rest frames (elementos estaticos)",
    "Vergence-accommodation: distancia fixa",
    "Comfort ratings: Meta/Oculus classifica",
    "  apps de Comfortable a Intense"
]
for item in items:
    add_para(tf, "\u2022 " + item, 13, LIGHT_GRAY, space_before=Pt(4))

add_source(slide, "Fontes: Reason & Brand, 1975; Kennedy SSQ; LaViola, 2000 (Cybersickness in VR); Frontiers in VR 2024 (1478106)")


# ── Requisitos tecnicos VR ──
slide = content_slide("Requisitos Tecnicos Minimos para VR",
    "Explicar cada parametro. Usar analogias medicas quando possivel.")

add_table(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.2),
    9, 4,
    ["Parametro", "Minimo", "Recomendado", "Porque importa"],
    [
        ["Frame rate", "60 fps", "90-120 fps", "Abaixo: nausea (flicker percebido)"],
        ["Latencia (MTP)", "< 20 ms", "< 12 ms", "Atraso entre mover cabeca e imagem atualizar"],
        ["Resolucao (por olho)", "1440x1600", "2064x2208+", "Reduz screen-door effect"],
        ["Field of View (FOV)", "90 graus", "110+ graus", "FOV estreito = visao tunel, quebra imersao"],
        ["Refresh rate", "72 Hz", "90-120 Hz", "Hz = quantas vezes o ecra atualiza por segundo"],
        ["Tracking", "3-DOF (rotacao)", "6-DOF (rot.+pos.)", "6-DOF permite andar; 3-DOF so olhar"],
        ["IPD adjustment", "Fixo", "Ajustavel", "Distancia inter-pupilar: errado = fadiga ocular"],
        ["Audio", "Stereo", "Espacial HRTF", "Audio 3D: som vem da direcao correta"],
    ])

tf = add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(1),
    "Glossario rapido:", 14, ACCENT, True)
add_para(tf, "fps = frames per second | MTP = motion-to-photon latency | DOF = degrees of freedom | "
         "FOV = field of view | IPD = inter-pupillary distance | HRTF = head-related transfer function", 12, SOURCE_CLR)


# ── Frame rate ──
slide = content_slide("Frame Rate: Cinema vs VR")
add_table(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(3.2),
    7, 3,
    ["Frame Rate", "Contexto", "Em VR"],
    [
        ["24 fps", "Cinema (standard desde 1927)", "INUTILIZAVEL — nausea severa"],
        ["30 fps", "TV / streaming video", "Nausea moderada a severa"],
        ["60 fps", "PC gaming / minimo VR", "Funcional mas desconfortavel"],
        ["72 fps", "Meta Quest 2 (default)", "Aceitavel para a maioria"],
        ["90 fps", "Meta Quest 3 / PCVR", "Bom conforto — objetivo da industria"],
        ["120 fps", "Quest 3 / PSVR2 / premium", "25% menos queixas vs 90fps (OLED)"],
    ])

tf = add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(2),
    "Porque 24fps funciona no cinema mas nao em VR?", 18, ACCENT, True)
add_para(tf, "No cinema, estamos FORA da imagem — o cerebro aceita 24fps num ecra plano.", 16, LIGHT_GRAY)
add_para(tf, "Em VR, estamos DENTRO — a cabeca roda, e o mundo deve responder instantaneamente.", 16, LIGHT_GRAY)
add_para(tf, "A 24fps, cada frame dura 41.7ms — o sistema vestibular deteta o atraso e entra em conflito.", 16, LIGHT_GRAY)
add_source(slide, "Fontes: Meta Performance Guidelines (developers.meta.com); IEEE TVCG 2023 (Frame Rate Effect on VR Sickness)")


# ── Escala em VR ──
slide = content_slide("Escala em VR: Porque 1:1 Importa",
    "A escala errada causa desconforto e quebra de imersao.")

tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(5.5),
    "O problema da escala", 20, ACCENT, True)
add_para(tf, "", 6, WHITE)
add_para(tf, "Em VR, 1 unidade = 1 metro (por convencao).", 17, LIGHT_GRAY)
add_para(tf, "", 6, WHITE)
add_para(tf, "Se um orgao estiver 2x maior que o real:", 17, WHITE, True)
add_para(tf, "\u2022 O cerebro estranha (uncanny valley espacial)", 16, LIGHT_GRAY)
add_para(tf, "\u2022 A aprendizagem anatomica fica comprometida", 16, LIGHT_GRAY)
add_para(tf, "\u2022 A percepcao de distancia falha", 16, LIGHT_GRAY)
add_para(tf, "", 8, WHITE)
add_para(tf, "Se uma sala estiver demasiado pequena:", 17, WHITE, True)
add_para(tf, "\u2022 Claustrofobia digital", 16, LIGHT_GRAY)
add_para(tf, "\u2022 Desconforto e ansiedade", 16, LIGHT_GRAY)
add_para(tf, "", 8, WHITE)
add_para(tf, "Best practice:", 17, ACCENT, True)
add_para(tf, "Medir objetos reais, replicar escala 1:1 em VR.", 16, LIGHT_GRAY)
add_para(tf, "Exemplo: coracao humano \u2248 12-14cm \u2192 scale=\"0.12\"", 16, ACCENT2)

add_card(slide, Inches(7.5), Inches(1.6), Inches(5.2), Inches(5.2), BG_CARD, ACCENT)
tf = add_textbox(slide, Inches(7.8), Inches(1.8), Inches(4.6), Inches(0.5),
    "Construir uma cena 3D", 18, ACCENT, True)
items = [
    "Sistema de coordenadas:",
    "  X = esquerda/direita",
    "  Y = baixo/cima",
    "  Z = frente/tras",
    "",
    "Unidades: metros (A-Frame)",
    "Olhos do utilizador: Y = 1.6m",
    "",
    "Boas praticas:",
    "  Objetos entre 1m e 10m de distancia",
    "  Texto legivel: >0.5m de largura",
    "  Nada demasiado perto (<0.5m)",
    "  Nada atras do utilizador no inicio",
    "  Iluminacao: min. 2 fontes de luz"
]
for item in items:
    add_para(tf, item, 14, LIGHT_GRAY, space_before=Pt(4))


# ── Gaming <-> Healthcare ──
slide = content_slide("Gaming e Saude: A Mesma Tecnologia")
pairs = [
    ("GPU de jogos AAA (NVIDIA, AMD)", "Renderiza simulacoes cirurgicas a 90fps", ACCENT),
    ("Audio espacial (Dolby Atmos, Wwise)", "Auscultacao virtual precisa", ACCENT2),
    ("Motion capture (Vicon, OptiTrack)", "Avaliacao de reabilitacao / marcha", ACCENT3),
    ("3D modeling (Blender, ZBrush)", "Modelos anatomicos detalhados", ACCENT4),
]
for i, (game, health, color) in enumerate(pairs):
    y = Inches(1.6 + i * 1.3)
    add_card(slide, Inches(0.8), y, Inches(5.2), Inches(1.1), BG_CARD, color)
    add_textbox(slide, Inches(1.1), y + Inches(0.15), Inches(4.6), Inches(0.7), game, 16, color, True)
    add_textbox(slide, Inches(6.1), y + Inches(0.15), Inches(0.8), Inches(0.7), "\u279C", 24, color, alignment=PP_ALIGN.CENTER)
    add_card(slide, Inches(7), y, Inches(5.5), Inches(1.1), BG_CARD, color)
    add_textbox(slide, Inches(7.3), y + Inches(0.15), Inches(5), Inches(0.7), health, 16, LIGHT_GRAY)

add_source(slide, "A industria gaming foi avaliada em $184B (2023, Newzoo). A mesma pipeline tecnologica alimenta VR em saude.")


# ── Demo ──
slide = content_slide("Demo ao Vivo")
tf = add_textbox(slide, Inches(2), Inches(2), Inches(9), Inches(4),
    "Momento de demonstracao", 36, ACCENT, True, PP_ALIGN.CENTER)
add_para(tf, "", 12, WHITE)
add_para(tf, "Mostrar a cena VR que vamos construir hoje", 22, WHITE, alignment=PP_ALIGN.CENTER)
add_para(tf, "Abrir no browser + entrar em modo VR no Quest", 22, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_para(tf, "", 12, WHITE)
add_para(tf, "Convidar 2-3 alunos a experimentar o headset", 20, ACCENT2, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  BLOCO 2 — FERRAMENTAS DE CONTEUDO
# ══════════════════════════════════════════════════════════════
section_slide("02", "Ferramentas de Criacao de Conteudo",
    "Som, imagem 360, video, 3D, Gaussian Splatting, Kiri Engine")

# ── Tipos de conteudo ──
slide = content_slide("Tipos de Conteudo VR")
types = [
    ("\U0001F50A", "Som", "Audio ambiente, espacial,\nHRTF, ambisonics", ACCENT),
    ("\U0001F5BC", "Imagem 360", "Panoramas equiretangulares,\nfotografia imersiva", ACCENT2),
    ("\U0001F3AC", "Video", "360, volumetrico,\nprocedimentos medicos", ACCENT3),
    ("\U0001F4A0", "3D", "CGI, fotogrametria,\nGaussian Splatting", ACCENT4),
]
for i, (icon, title, desc, color) in enumerate(types):
    x = Inches(0.8 + i * 3.05)
    add_card(slide, x, Inches(1.6), Inches(2.8), Inches(4), BG_CARD, color)
    add_textbox(slide, x + Inches(0.3), Inches(1.9), Inches(2.2), Inches(0.8), icon, 40, color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.3), Inches(2.8), Inches(2.2), Inches(0.5), title, 22, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.3), Inches(3.5), Inches(2.2), Inches(1.5), desc, 14, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ── Som ──
slide = content_slide("Som em VR: 50% da Experiencia Imersiva")
add_table(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(2.5),
    5, 3,
    ["Tipo (termo tecnico)", "Descricao", "Exemplo clinico"],
    [
        ["Mono/Stereo", "Audio standard, sem posicionamento espacial", "Musica de fundo"],
        ["Positional Audio (3D)", "Som posicionado no espaco 3D, atenua com distancia", "Batimento cardiaco vindo do peito virtual"],
        ["Ambisonics (B-format)", "Captura/reproducao esferica completa do campo sonoro", "Gravacao do ambiente de bloco operatorio"],
        ["HRTF (Head-Related TF)", "Modela a interacao do som com pinnae, cabeca e torso", "Distinguir sons acima vs abaixo do doente"],
    ])

tf = add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(2.5),
    "Aplicacoes medicas do audio espacial:", 18, ACCENT2, True)
add_para(tf, "\u2022 Auscultacao virtual — som mais alto quando estetoscopio no ponto correto", 16, LIGHT_GRAY)
add_para(tf, "\u2022 Simulacao de bloco — bips de monitores, ventilador, comunicacao (stress autentico)", 16, LIGHT_GRAY)
add_para(tf, "\u2022 Terapia: paisagens sonoras 360 para gestao de dor e ansiedade", 16, LIGHT_GRAY)
add_para(tf, "\u2022 Reabilitacao: instrucoes auditivas espaciais ('alcance o som a sua direita')", 16, LIGHT_GRAY)
add_source(slide, "Fonte: Audio Engineering Society; Frontiers in VR 2025 (1629908) — spatial audio significantly improves immersion")


# ── Imagem 360 ──
slide = content_slide("Imagem 360: Projecao Equiretangular")
tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(7), Inches(5),
    "Projecao Equiretangular (Equirectangular Projection)", 18, ACCENT, True)
add_para(tf, "", 6, WHITE)
add_para(tf, "Analogia: descascar um globo e esticar num retangulo (mapa-mundi).", 16, LIGHT_GRAY)
add_para(tf, "A imagem fica distorcida em 2D, mas quando envolvida numa", 16, LIGHT_GRAY)
add_para(tf, "esfera em VR (<a-sky>), parece perfeitamente natural.", 16, LIGHT_GRAY)
add_para(tf, "", 8, WHITE)
add_para(tf, "Especificacoes tecnicas:", 16, WHITE, True)
add_para(tf, "\u2022 Racio: 2:1 (ex: 4096x2048 ou 8192x4096 pixels)", 15, LIGHT_GRAY)
add_para(tf, "\u2022 Formato: JPEG ou PNG (JPEG mais leve)", 15, LIGHT_GRAY)
add_para(tf, "\u2022 Projecao: latitude-longitude mapping", 15, LIGHT_GRAY)
add_para(tf, "", 8, WHITE)
add_para(tf, "Equipamento acessivel:", 16, WHITE, True)
add_para(tf, "\u2022 Insta360 X3/X4 (~400-500 EUR)", 15, LIGHT_GRAY)
add_para(tf, "\u2022 Ricoh Theta SC2 (~300 EUR)", 15, LIGHT_GRAY)
add_para(tf, "\u2022 GoPro MAX (~400 EUR)", 15, LIGHT_GRAY)
add_para(tf, "\u2022 Smartphone apps: Google Street View, Panorama 360", 15, LIGHT_GRAY)

add_card(slide, Inches(8.2), Inches(1.6), Inches(4.5), Inches(5), BG_CARD, ACCENT2)
tf = add_textbox(slide, Inches(8.5), Inches(1.8), Inches(4), Inches(0.5),
    "Aplicacoes em Saude", 18, ACCENT2, True)
items = [
    "Visitas virtuais ao hospital",
    "Documentacao do bloco operatorio",
    "Ambientes clinicos para estagios",
    "Telemedicina: ambiente do doente",
    "Documentacao forense 360"
]
for item in items:
    add_para(tf, "\u2022 " + item, 14, LIGHT_GRAY, space_before=Pt(8))


# ── Video ──
slide = content_slide("Video 360 vs Video Volumetrico")
add_table(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.5),
    8, 3,
    ["Parametro", "Video 360", "Video Volumetrico"],
    [
        ["Custo producao", "Baixo-Medio (camera 360)", "Alto (multi-camera studio)"],
        ["Equipamento", "300-5.000 EUR (Insta360, Theta)", "100K+ EUR (Azure Kinect array)"],
        ["Graus de liberdade", "3-DOF (olhar a volta)", "6-DOF (mover-se a volta)"],
        ["Tamanho ficheiro", "Grande (gerivel)", "Muito grande"],
        ["Realismo", "Alto (filmagem real)", "Maximo (filmagem real 3D)"],
        ["Interatividade", "Limitada (hotspots)", "Alta (perspetiva livre)"],
        ["Melhor para", "Observacao procedimentos, tours", "Treino cirurgico, exame fisico"],
    ])
add_source(slide, "Exemplo: projeto THRIVE — video volumetrico de profissionais medicos para treino (AIXR, 2024)")


# ── Pipeline 3D ──
slide = content_slide("Pipeline de Criacao 3D (CG Pipeline)")
steps = [
    ("1", "Modeling", "Criar geometria\n(poligonos, NURBS,\nescultura digital)", ACCENT),
    ("2", "Texturing", "UV mapping,\nmateriais PBR\n(albedo, normal, AO)", ACCENT2),
    ("3", "Rigging", "Armature/skeleton\npara deformacao\nde mesh", ACCENT3),
    ("4", "Animation", "Keyframes,\nmotion capture,\nprocedural", ACCENT4),
    ("5", "Lighting", "Luzes (point,\ndirectional, area,\nambient, HDRI)", ACCENT),
    ("6", "Rendering", "Rasterization\n(tempo real) ou\nray tracing", ACCENT2),
]
for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 2.1)
    add_card(slide, x, Inches(1.6), Inches(1.9), Inches(4), BG_CARD, color)
    add_textbox(slide, x + Inches(0.2), Inches(1.8), Inches(1.5), Inches(0.6), num, 36, color, True, PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(2.5), Inches(1.5), Inches(0.5), title, 16, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(3.1), Inches(1.5), Inches(2), desc, 12, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    if i < 5:
        add_textbox(slide, x + Inches(1.85), Inches(2.8), Inches(0.5), Inches(0.5), "\u279C", 20, MID_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(5.9), Inches(11), Inches(0.5),
    "PBR = Physically Based Rendering | NURBS = Non-Uniform Rational B-Splines | AO = Ambient Occlusion | UV = coordenadas de textura",
    11, SOURCE_CLR)


# ── Maya vs Blender ──
slide = content_slide("Maya vs Blender: Comparacao")
add_table(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4),
    7, 3,
    ["Parametro", "Autodesk Maya", "Blender"],
    [
        ["Preco", "~2.016 EUR/ano (subscription)", "GRATUITO (GPL, para sempre)"],
        ["Curva de aprendizagem", "Ingreme, interface complexa", "Moderada (UI renovado desde 2.8)"],
        ["Uso industria", "Standard: Disney, Pixar, ILM", "Crescendo: Ubisoft, Netflix, NASA"],
        ["Pontos fortes", "NURBS, rigging, interop Autodesk", "Tudo-em-um: model, sculpt, video, VFX"],
        ["Linguagem script", "MEL / Python", "Python (mesmo!)"],
        ["Saude: recomendacao", "Grandes instituicoes com licencas", "IDEAL para medicos e investigadores"],
    ])
add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.6),
    "Recomendacao para esta aula: Blender (gratuito, acessivel, comunidade enorme, tutoriais em PT)", 16, ACCENT2, True)
add_source(slide, "Blender Foundation (blender.org) | Autodesk Maya (autodesk.com/maya)")


# ── Gaussian Splatting — historia e ciencia ──
slide = content_slide("Gaussian Splatting: Historia e Ciencia",
    "Explicar a evolucao desde NeRF ate 3DGS. Os medicos apreciam a timeline.")

tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(6.2), Inches(5.5),
    "Evolucao historica", 20, ACCENT, True)
add_para(tf, "", 4, WHITE)
add_para(tf, "2020 — NeRF (Neural Radiance Fields)", 16, WHITE, True)
add_para(tf, "  Mildenhall et al. (UC Berkeley). Artigo seminal em ECCV 2020.", 14, LIGHT_GRAY)
add_para(tf, "  Rede neural que aprende a representar cenas 3D a partir de fotos.", 14, LIGHT_GRAY)
add_para(tf, "  Problema: rendering MUITO lento (minutos por frame).", 14, LIGHT_GRAY)
add_para(tf, "", 6, WHITE)
add_para(tf, "2021-2022 — Instant-NGP, Plenoxels, TensoRF", 16, WHITE, True)
add_para(tf, "  Otimizacoes (NVIDIA Instant-NGP: treino em segundos).", 14, LIGHT_GRAY)
add_para(tf, "  Ainda baseados em redes neurais \u2192 rendering lento.", 14, LIGHT_GRAY)
add_para(tf, "", 6, WHITE)
add_para(tf, "2023 — 3D Gaussian Splatting (3DGS)", 16, ACCENT, True)
add_para(tf, "  Kerbl, Kopanas, Leimkuhler & Drettakis", 14, WHITE, True)
add_para(tf, "  INRIA (Universite Cote d'Azur, Franca). SIGGRAPH 2023.", 14, LIGHT_GRAY)
add_para(tf, "  Revolucao: SEM rede neural. Usa milhoes de Gaussianas 3D", 14, LIGHT_GRAY)
add_para(tf, "  (elipsoides semi-transparentes) = rendering TEMPO REAL.", 14, LIGHT_GRAY)
add_para(tf, "  138 fps em cenas complexas vs <1 fps com NeRF.", 14, ACCENT2, True)
add_para(tf, "", 6, WHITE)
add_para(tf, "2024-2025 — Adocao massiva", 16, WHITE, True)
add_para(tf, "  Apps mobile (Kiri, Polycam), medicina (MedGS),", 14, LIGHT_GRAY)
add_para(tf, "  Google Maps, Apple Vision Pro, automacao industrial.", 14, LIGHT_GRAY)

add_card(slide, Inches(7.5), Inches(1.6), Inches(5.2), Inches(5.2), BG_CARD, ACCENT)
tf = add_textbox(slide, Inches(7.8), Inches(1.8), Inches(4.6), Inches(0.5),
    "O que e uma Gaussiana 3D?", 16, ACCENT, True)
add_para(tf, "", 4, WHITE)
add_para(tf, "Matematicamente: funcao de distribuicao", 14, WHITE, True)
add_para(tf, "normal tridimensional definida por:", 14, WHITE, True)
add_para(tf, "", 4, WHITE)
add_para(tf, "\u2022 Posicao (x, y, z) — centro no espaco", 13, LIGHT_GRAY)
add_para(tf, "\u2022 Covariancia (3x3 matrix) — forma/orientacao", 13, LIGHT_GRAY)
add_para(tf, "\u2022 Opacidade (alpha) — transparencia", 13, LIGHT_GRAY)
add_para(tf, "\u2022 Cor (spherical harmonics) — cor dependente da direcao de vista", 13, LIGHT_GRAY)
add_para(tf, "", 6, WHITE)
add_para(tf, "Rendering: projecao das Gaussianas", 14, WHITE, True)
add_para(tf, "3D em splats 2D no ecra (splatting).", 14, WHITE, True)
add_para(tf, "Diferenciavel \u2192 otimizavel por", 14, LIGHT_GRAY)
add_para(tf, "gradient descent.", 14, LIGHT_GRAY)
add_para(tf, "", 6, WHITE)
add_para(tf, "Vantagem chave:", 14, ACCENT2, True)
add_para(tf, "Rasterization-based (GPU-friendly)", 14, LIGHT_GRAY)
add_para(tf, "vs ray marching (NeRF, lento).", 14, LIGHT_GRAY)

add_source(slide, "Kerbl et al., '3D Gaussian Splatting for Real-Time Radiance Field Rendering', SIGGRAPH 2023 (repo.inria.fr/fungraph/3d-gaussian-splatting)")


# ── Gaussian Splatting aplicacoes ──
slide = content_slide("Gaussian Splatting: Aplicacoes Conhecidas")
items_gs = [
    ("Saude / Medicina", "MedGS: reconstrucao de ecografia, RM, TC em 3D\nAnatomia cinematica em dispositivos moveis\nImageologia dentaria de alta fidelidade (Jin et al., 2025)", ACCENT),
    ("Google Maps / Earth", "Immersive View: reconstrucao de cidades inteiras\ncom Gaussian Splatting a partir de imagens Street View", ACCENT2),
    ("Imobiliario / Arquitetura", "Visitas virtuais fotorrealistas de imoveis e edificios\nMatterport, Zillow ja integram GS", ACCENT3),
    ("Cinema / VFX", "Captura de ambientes reais para pos-producao\nIntegracao com pipelines Unreal Engine 5", ACCENT4),
]
for i, (title, desc, color) in enumerate(items_gs):
    y = Inches(1.6 + i * 1.4)
    add_card(slide, Inches(0.8), y, Inches(11.5), Inches(1.2), BG_CARD, color)
    add_textbox(slide, Inches(1.1), y + Inches(0.1), Inches(3.5), Inches(1), title, 18, color, True)
    add_textbox(slide, Inches(4.8), y + Inches(0.1), Inches(7.2), Inches(1), desc, 14, LIGHT_GRAY)

add_source(slide, "Fontes: SAGE 2025 (dental GS); MedGS (arXiv 2509.16806); Google Immersive View (blog.google)")


# ── Kiri Engine ──
slide = content_slide("Kiri Engine: Scanner 3D no Telemovel")
tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(5),
    "O que e?", 20, ACCENT, True)
add_para(tf, "App gratuita (Android/iOS) — scanner 3D com o telemovel.", 16, LIGHT_GRAY, space_before=Pt(10))
add_para(tf, "", 6, WHITE)
add_para(tf, "Tecnologias suportadas:", 16, WHITE, True)
add_para(tf, "\u2022 Fotogrametria — qualquer telemovel com camera", 15, LIGHT_GRAY)
add_para(tf, "\u2022 LiDAR scanning — iPhone 12 Pro+ (mais rapido)", 15, LIGHT_GRAY)
add_para(tf, "\u2022 3D Gaussian Splatting — primeiro GS mobile do mundo", 15, LIGHT_GRAY)
add_para(tf, "\u2022 Neural Surface Reconstruction — AI-enhanced", 15, LIGHT_GRAY)
add_para(tf, "", 6, WHITE)
add_para(tf, "Processamento: ~2 min/scan | Exporta: OBJ, GLB, USDZ", 15, ACCENT2, True)
add_para(tf, "", 6, WHITE)
add_para(tf, "Precos:", 16, WHITE, True)
add_para(tf, "\u2022 Gratis: fotogrametria + LiDAR, 3 exports/semana", 15, LIGHT_GRAY)
add_para(tf, "\u2022 Pro: 17.99 EUR/mes — GS + ilimitado", 15, LIGHT_GRAY)

add_card(slide, Inches(7.5), Inches(1.6), Inches(5.2), Inches(5), BG_CARD, ACCENT2)
tf = add_textbox(slide, Inches(7.8), Inches(1.8), Inches(4.6), Inches(0.5),
    "Alternativas:", 18, ACCENT2, True)
items = [
    "Polycam — iOS/Android, LiDAR, GS",
    "Luma AI — Gaussian Splatting online",
    "Meshroom (AliceVision) — gratuito, PC",
    "RealityScan (Epic Games) — fotogrametria",
    "3D Scanner App (Apple) — LiDAR nativo",
    "Qlone — scan com QR code de referencia",
    "Regard3D — open-source, PC",
]
for item in items:
    add_para(tf, "\u2022 " + item, 14, LIGHT_GRAY, space_before=Pt(6))

add_source(slide, "Kiri Engine (kiriengine.app) | HTC VIVERSE partnership, Dec 2025 (businesswire.com)")


# ── Exercicio Kiri ──
slide = content_slide("EXERCICIO: Scan 3D da Sala com Kiri")
tf = add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(5),
    "Instrucoes", 28, ACCENT, True, PP_ALIGN.CENTER)
add_para(tf, "", 12, WHITE)
add_para(tf, "1. Formem grupos de 3-4 pessoas", 22, WHITE, alignment=PP_ALIGN.CENTER)
add_para(tf, "2. Instalem Kiri Engine no telemovel (1 por grupo)", 22, WHITE, alignment=PP_ALIGN.CENTER)
add_para(tf, "3. Escolham um objeto ou area da sala", 22, WHITE, alignment=PP_ALIGN.CENTER)
add_para(tf, "4. Tirem 30-50 fotos de angulos diferentes", 22, WHITE, alignment=PP_ALIGN.CENTER)
add_para(tf, "5. Upload e aguardem (~2 min)", 22, WHITE, alignment=PP_ALIGN.CENTER)
add_para(tf, "", 12, WHITE)
add_para(tf, "Este scan sera o ambiente VR da proxima aula!", 20, ACCENT2, True, alignment=PP_ALIGN.CENTER)
add_para(tf, "Tempo: 15 minutos", 20, ACCENT3, True, alignment=PP_ALIGN.CENTER)


# ── PAUSA ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_textbox(slide, Inches(2), Inches(2.5), Inches(9), Inches(2), "PAUSA", 72, ACCENT, True, PP_ALIGN.CENTER, "Segoe UI Light")
add_textbox(slide, Inches(2), Inches(4.5), Inches(9), Inches(1), "10 minutos", 28, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  BLOCO 3 — INTERATIVIDADE
# ══════════════════════════════════════════════════════════════
section_slide("03", "Interatividade e Motores de Jogo",
    "Game Engines, Unity, Unreal, A-Frame, Inspector")


# ── Game Engine definicao ──
slide = content_slide("O que e um Game Engine (Motor de Jogo)?",
    "Definicao tecnica. Estes motores sao a base de toda VR interativa.")

tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(1.5),
    "Definicao", 22, ACCENT, True)
add_para(tf, "Um Game Engine e um framework de software que fornece as funcionalidades", 17, LIGHT_GRAY)
add_para(tf, "base para criar aplicacoes interativas em tempo real:", 17, LIGHT_GRAY)

add_card(slide, Inches(0.8), Inches(3.4), Inches(3.5), Inches(3.2), BG_CARD, ACCENT)
tf = add_textbox(slide, Inches(1.1), Inches(3.6), Inches(3), Inches(0.4), "Rendering Engine", 16, ACCENT, True)
add_para(tf, "Desenha graficos 2D/3D em tempo real (rasterization ou ray tracing)", 13, LIGHT_GRAY, space_before=Pt(6))

add_card(slide, Inches(4.6), Inches(3.4), Inches(3.5), Inches(3.2), BG_CARD, ACCENT2)
tf = add_textbox(slide, Inches(4.9), Inches(3.6), Inches(3), Inches(0.4), "Physics Engine", 16, ACCENT2, True)
add_para(tf, "Simula gravidade, colisoes, corpos rigidos e fluidos", 13, LIGHT_GRAY, space_before=Pt(6))

add_card(slide, Inches(8.4), Inches(3.4), Inches(4.2), Inches(3.2), BG_CARD, ACCENT3)
tf = add_textbox(slide, Inches(8.7), Inches(3.6), Inches(3.7), Inches(0.4), "Mais subsistemas:", 16, ACCENT3, True)
items = ["Audio engine (espacial)", "Input system (teclado, VR controllers)",
         "Networking (multiplayer)", "Scripting (logica do jogo)",
         "Asset pipeline (modelos, texturas)", "Scene graph (hierarquia de objetos)"]
for item in items:
    add_para(tf, "\u2022 " + item, 12, LIGHT_GRAY, space_before=Pt(3))

add_source(slide, "Definicao adaptada de Gregory, J. 'Game Engine Architecture' (3rd ed., CRC Press, 2018)")


# ── Unity ──
slide = content_slide("Unity (C#)")
add_card(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5), BG_CARD, ACCENT)
tf = add_textbox(slide, Inches(1.1), Inches(1.8), Inches(5), Inches(0.5), "Caracteristicas", 20, ACCENT, True)
items = ["Motor mais usado para VR no mundo", "C# (C-Sharp) — linguagem orientada a objetos",
         "Maior ecossistema: Asset Store com milhares de assets", "Ideal para Meta Quest (standalone)",
         "Cross-platform: PC, mobile, consolas, web", "Versao gratuita para projetos < $200K/ano"]
for item in items:
    add_para(tf, "\u2022 " + item, 14, LIGHT_GRAY, space_before=Pt(5))

add_card(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(5), BG_CARD, ACCENT)
tf = add_textbox(slide, Inches(7.1), Inches(1.8), Inches(5), Inches(0.5), "Saude: casos de uso", 20, ACCENT, True)
items = ["Osso VR — treino cirurgico ortopedico", "SimX — simulacao de emergencias medicas",
         "ImmersiveTouch — neurocirurgia VR", "Apps AR anatomia (Complete Anatomy)",
         "Jogos de reabilitacao gamificada", "Simulacao de dispositivos medicos"]
for item in items:
    add_para(tf, "\u2022 " + item, 14, LIGHT_GRAY, space_before=Pt(5))
add_source(slide, "Unity Technologies (unity.com) | Versao LTS recomendada para producao")


# ── Unreal ──
slide = content_slide("Unreal Engine (C++ / Blueprints)")
add_card(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5), BG_CARD, ACCENT3)
tf = add_textbox(slide, Inches(1.1), Inches(1.8), Inches(5), Inches(0.5), "Caracteristicas", 20, ACCENT3, True)
items = ["Qualidade grafica SUPERIOR (fotorrealismo)", "C++ (performance maxima) + Blueprints (visual scripting)",
         "Nanite: geometria virtualizada (trilioes de poligonos)", "Lumen: iluminacao global em tempo real",
         "MetaHuman: criacao de humanos fotorrealistas", "Gratuito ate $1M receita (5% royalty depois)"]
for item in items:
    add_para(tf, "\u2022 " + item, 14, LIGHT_GRAY, space_before=Pt(5))

add_card(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(5), BG_CARD, ACCENT3)
tf = add_textbox(slide, Inches(7.1), Inches(1.8), Inches(5), Inches(0.5), "Saude: casos de uso", 20, ACCENT3, True)
items = ["Simulacao cirurgica alta fidelidade", "Visualizacao anatomica fotorrealista",
         "Digital twins de hospitais/blocos", "Encontros virtuais com doentes realistas (MetaHuman)",
         "Simulacao de cenarios de emergencia", "Requer hardware potente (GPU dedicada)"]
for item in items:
    add_para(tf, "\u2022 " + item, 14, LIGHT_GRAY, space_before=Pt(5))
add_source(slide, "Epic Games (unrealengine.com) | Blueprints = visual scripting SEM codigo")


# ── Tabela comparativa ──
slide = content_slide("Comparacao: Unity vs Unreal vs A-Frame")
add_table(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.8),
    9, 4,
    ["", "Unity", "Unreal Engine", "A-Frame"],
    [
        ["Linguagem", "C#", "C++ / Blueprints", "HTML / JavaScript"],
        ["Instalacao", "Sim (2-10 GB)", "Sim (30-50 GB)", "NAO (browser)"],
        ["Curva aprendizagem", "Media", "Dificil", "FACIL"],
        ["Qualidade grafica", "Boa", "Excelente", "Basica a Boa"],
        ["Suporte VR", "Excelente", "Excelente", "Bom (WebXR)"],
        ["Preco", "Gratis (< $200K)", "Gratis (< $1M)", "GRATUITO (sempre)"],
        ["Melhor para", "Apps producao, Quest", "Simulacao high-end", "APRENDER, prototipar"],
        ["Tempo ate 1.a cena", "Horas", "Horas", "MINUTOS"],
    ])


# ── O que e HTML ──
slide = content_slide("HTML: A Linguagem da Web",
    "Explicar de raiz para quem nunca programou.")

tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(5.5),
    "HyperText Markup Language", 20, ACCENT, True)
add_para(tf, "", 4, WHITE)
add_para(tf, "NAO e uma linguagem de programacao.", 17, ACCENT3, True)
add_para(tf, "E uma linguagem de MARCACAO — descreve a ESTRUTURA", 17, LIGHT_GRAY)
add_para(tf, "de um documento (como XML, como um formulario clinico).", 17, LIGHT_GRAY)
add_para(tf, "", 8, WHITE)
add_para(tf, "Conceitos fundamentais:", 17, WHITE, True)
add_para(tf, "", 4, WHITE)
add_para(tf, "TAG — instrucao entre < >", 16, ACCENT2, True)
add_para(tf, '  <title>Aula VR</title>', 14, ACCENT2)
add_para(tf, '  ^abertura  ^conteudo  ^fecho', 13, MID_GRAY)
add_para(tf, "", 6, WHITE)
add_para(tf, "ATRIBUTO — propriedade dentro da tag", 16, ACCENT2, True)
add_para(tf, '  <a-box color="red" position="0 1 -5">', 14, ACCENT2)
add_para(tf, '         ^nome="valor"  ^nome="valor"', 13, MID_GRAY)
add_para(tf, "", 6, WHITE)
add_para(tf, "NESTING — tags dentro de tags (hierarquia)", 16, ACCENT2, True)
add_para(tf, '  <body>', 14, ACCENT2)
add_para(tf, '    <a-scene>', 14, ACCENT2)
add_para(tf, '      <a-box></a-box>  <- filho de scene', 14, ACCENT2)
add_para(tf, '    </a-scene>', 14, ACCENT2)
add_para(tf, '  </body>', 14, ACCENT2)

add_card(slide, Inches(7.5), Inches(1.6), Inches(5.2), Inches(5.2), BG_CARD, ACCENT)
tf = add_textbox(slide, Inches(7.8), Inches(1.8), Inches(4.6), Inches(0.5),
    "HTML + CSS + JS", 18, ACCENT, True)
add_para(tf, "", 6, WHITE)
add_para(tf, "HTML = Estrutura (esqueleto)", 16, ACCENT, True)
add_para(tf, "O QUE existe: tags, conteudo", 14, LIGHT_GRAY)
add_para(tf, "Analogia: campos do processo clinico", 14, MID_GRAY, italic=True)
add_para(tf, "", 8, WHITE)
add_para(tf, "CSS = Aparencia (pele)", 16, ACCENT2, True)
add_para(tf, "COMO se ve: cores, tamanhos", 14, LIGHT_GRAY)
add_para(tf, "Analogia: formatacao do relatorio", 14, MID_GRAY, italic=True)
add_para(tf, "", 8, WHITE)
add_para(tf, "JavaScript = Comportamento (musculos)", 16, ACCENT3, True)
add_para(tf, "O QUE FAZ: logica, interacoes", 14, LIGHT_GRAY)
add_para(tf, "Analogia: instrucoes de tratamento", 14, MID_GRAY, italic=True)
add_para(tf, "", 10, WHITE)
add_para(tf, "Nesta aula: usamos so HTML!", 16, ACCENT, True)
add_para(tf, "A-Frame transforma HTML em 3D/VR.", 14, LIGHT_GRAY)


# ── A-Frame ──
slide = content_slide("A-Frame: VR com HTML")
tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(5),
    "Framework open-source para WebXR", 18, ACCENT, True)
add_para(tf, "Criado por Mozilla (2015), mantido pela comunidade.", 15, LIGHT_GRAY, space_before=Pt(8))
add_para(tf, "Construido sobre Three.js (biblioteca 3D JavaScript).", 15, LIGHT_GRAY)
add_para(tf, "", 8, WHITE)
add_para(tf, "Exemplo — cubo vermelho em VR:", 16, WHITE, True)
add_para(tf, '<a-box position="0 1 -5" color="red"></a-box>', 16, ACCENT2, space_before=Pt(6))
add_para(tf, "Sim, e so isto. Uma linha de HTML = um objeto 3D.", 16, ACCENT, True, space_before=Pt(8))
add_para(tf, "", 8, WHITE)
add_para(tf, "Suporte WebXR (W3C standard):", 16, WHITE, True)
add_para(tf, "\u2022 Chrome, Edge, Firefox, Safari (2025+)", 15, LIGHT_GRAY)
add_para(tf, "\u2022 Meta Quest Browser (nativo)", 15, LIGHT_GRAY)
add_para(tf, "\u2022 Funciona em PC, mobile E headset VR", 15, LIGHT_GRAY)

add_card(slide, Inches(7.5), Inches(1.6), Inches(5.2), Inches(5), BG_CARD, ACCENT)
tf = add_textbox(slide, Inches(7.8), Inches(1.8), Inches(4.6), Inches(0.5),
    "Vantagens:", 18, ACCENT, True)
items = ["\u2713 Baseado em HTML (familiar)", "\u2713 Zero instalacao", "\u2713 Partilha com URL",
         "\u2713 VR-ready automaticamente", "\u2713 WebXR completo (W3C)", "\u2713 Centenas de componentes comunitarios",
         "\u2713 Completamente gratuito (MIT license)", "\u2713 Resultado em MINUTOS"]
for item in items:
    add_para(tf, item, 15, LIGHT_GRAY, space_before=Pt(5))

# QR code
qr_path = make_qr("https://aframe.io")
slide.shapes.add_picture(qr_path, Inches(8.5), Inches(5.5), Inches(1.5), Inches(1.5))
add_textbox(slide, Inches(10.2), Inches(5.8), Inches(2.5), Inches(0.8),
    "aframe.io", 14, ACCENT, True)
add_source(slide, "A-Frame (aframe.io) | MIT License | Three.js (threejs.org) | WebXR Device API (W3C)")


# ── Glossario A-Frame ──
slide = content_slide("A-Frame: Glossario Tecnico (ECS Architecture)",
    "Entity-Component-System: a arquitetura base do A-Frame.")

add_table(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2),
    11, 3,
    ["Termo", "Definicao", "Exemplo"],
    [
        ["Entity (<a-entity>)", "Objeto generico/container. Sozinho nao faz nada.", "<a-entity></a-entity>"],
        ["Component", "Atributo que da funcionalidade a uma entity.", 'position="0 1 -5"'],
        ["Primitive", "Entity pre-configurada com componentes default.", "<a-box>, <a-sphere>, <a-sky>"],
        ["Scene (<a-scene>)", "Raiz: contem TUDO. Inicia o renderer WebGL.", "<a-scene>...</a-scene>"],
        ["Assets (<a-assets>)", "Pre-carrega ficheiros (imagens, modelos, audio).", '<img id="x" src="foto.jpg">'],
        ["Attribute", "Par nome=valor que configura um componente.", 'color="#FF0000"'],
        ["Animation", "Anima qualquer propriedade ao longo do tempo.", 'animation="property: rotation; ..."'],
        ["Event", "Acao que dispara algo (mouseenter, click, fusing).", '_event: mouseenter'],
        ["Mixin (<a-mixin>)", "Template reutilizavel de componentes.", '<a-mixin id="red" color="red">'],
        ["System", "Logica global singleton. Gere componentes do mesmo tipo.", "renderer system, physics system"],
    ])

add_source(slide, "Documentacao: aframe.io/docs | ECS = Entity-Component-System (padrao de arquitetura de game engines)")


# ── A-Frame: Primitivos e Tipos ──
slide = content_slide("A-Frame: Primitivos Disponiveis (Built-in)")

add_table(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5),
    12, 3,
    ["Primitivo", "Componentes implicitios", "Uso tipico"],
    [
        ["<a-box>", "geometry(box), material", "Edificios, paredes, caixas"],
        ["<a-sphere>", "geometry(sphere), material", "Planetas, bolas, atomos"],
        ["<a-cylinder>", "geometry(cylinder), material", "Pilares, tubos, vasos"],
        ["<a-plane>", "geometry(plane), material", "Chao, paredes, ecras"],
        ["<a-sky>", "geometry(sphere r:5000), material", "Panorama 360 (equirectangular)"],
        ["<a-torus>", "geometry(torus), material", "Aneis, donuts, orbitas"],
        ["<a-cone>", "geometry(cone), material", "Setas, indicadores"],
        ["<a-circle>", "geometry(circle), material", "Circulos, relogios"],
        ["<a-text>", "text", "Labels, titulos, instrucoes"],
        ["<a-image>", "geometry(plane), material(src)", "Fotos, diagramas 2D em VR"],
        ["<a-video>", "geometry(plane), material(src)", "Videos 2D dentro da cena VR"],
    ])


# ── Inspector ──
slide = content_slide("A-Frame Inspector: Editor Visual Integrado")

tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(5.5),
    "Ctrl + Alt + I", 36, ACCENT, True)
add_para(tf, "", 6, WHITE)
add_para(tf, "O Inspector e um editor visual BUILT-IN.", 17, LIGHT_GRAY)
add_para(tf, "Nao precisa instalar nada — esta incluido em", 17, LIGHT_GRAY)
add_para(tf, "QUALQUER cena A-Frame do mundo!", 17, LIGHT_GRAY)
add_para(tf, "", 10, WHITE)
add_para(tf, "Funcionalidades:", 17, WHITE, True)
add_para(tf, "\u2022 Scene graph: arvore hierarquica de entities", 15, LIGHT_GRAY)
add_para(tf, "\u2022 Component editor: editar propriedades com UI", 15, LIGHT_GRAY)
add_para(tf, "\u2022 Transform gizmos: mover/rodar/escalar visual", 15, LIGHT_GRAY)
add_para(tf, "\u2022 Viewport 3D: navegar a cena livremente", 15, LIGHT_GRAY)
add_para(tf, "\u2022 Copy to clipboard: copiar HTML do objeto", 15, LIGHT_GRAY)
add_para(tf, "", 10, WHITE)
add_para(tf, "Ideal para quem tem receio de 'partir o codigo'!", 17, ACCENT2, True)
add_para(tf, "Explorem visualmente, depois copiem valores.", 16, LIGHT_GRAY)

add_card(slide, Inches(7.5), Inches(1.6), Inches(5.2), Inches(5.2), BG_CARD, ACCENT2)
tf = add_textbox(slide, Inches(7.8), Inches(1.8), Inches(4.6), Inches(0.5),
    "Passo a passo:", 18, ACCENT2, True)
steps = [
    "1. Abrir a cena no Chrome",
    "2. Carregar Ctrl + Alt + I",
    "3. Painel esquerdo: lista de entities",
    "4. Clicar numa entity",
    "5. Painel direito: propriedades",
    "   (position, rotation, scale, color...)",
    "6. Editar valores com sliders/inputs",
    "7. Arrastar gizmos na viewport",
    "8. Ver alteracoes em tempo real",
    "9. Copiar valores para o HTML",
    "",
    "Atalhos no Inspector:",
    "  W = translate | E = rotate",
    "  R = scale | Delete = apagar",
]
for step in steps:
    add_para(tf, step, 14, LIGHT_GRAY, space_before=Pt(4))


# ══════════════════════════════════════════════════════════════
#  BLOCO 4 — EXERCICIO PRATICO
# ══════════════════════════════════════════════════════════════
section_slide("04", "Exercicio Pratico: A-Frame",
    "Vamos construir uma cena VR passo a passo — sem medo!")

# ── Abordagem ──
slide = content_slide("A Nossa Abordagem Pedagogica")
tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5.5),
    "Nao vao ser programadores.", 28, ACCENT3, True)
add_para(tf, "Vao COMPREENDER a ferramenta e ver resultados imediatos.", 24, ACCENT, True)
add_para(tf, "", 10, WHITE)
add_para(tf, "Metodologia (modify-existing-code, evidence-based):", 18, WHITE, True)
add_para(tf, "\u2022 Comecar com codigo que JA FUNCIONA (motivacao imediata)", 17, LIGHT_GRAY)
add_para(tf, "\u2022 Mudar UM valor de cada vez e observar o resultado", 17, LIGHT_GRAY)
add_para(tf, "\u2022 Copiar-colar padroes existentes e adaptar (scaffolding)", 17, LIGHT_GRAY)
add_para(tf, "\u2022 Usar o Inspector para explorar visualmente (Ctrl+Alt+I)", 17, LIGHT_GRAY)
add_para(tf, "", 10, WHITE)
add_para(tf, "Regras de ouro:", 18, WHITE, True)
add_para(tf, "\u2022 Ctrl+Z = desfazer (funciona SEMPRE)", 17, ACCENT2)
add_para(tf, "\u2022 Ctrl+S = guardar (ve resultado no browser)", 17, ACCENT2)
add_para(tf, "\u2022 Nada explode! Podem experimentar a vontade.", 17, ACCENT2)
add_source(slide, "Wilson, G. 'Ten Quick Tips for Teaching Programming', PLOS Computational Biology, 2019")


# ── Setup ──
slide = content_slide("Setup: VS Code + Live Server")
tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
    "Passos:", 22, ACCENT, True)
add_para(tf, "", 6, WHITE)
add_para(tf, "1. Abrir VS Code", 20, WHITE)
add_para(tf, "2. File > Open Folder > pasta do projeto", 20, WHITE)
add_para(tf, "3. Abrir o ficheiro index.html", 20, WHITE)
add_para(tf, "4. Clicar 'Go Live' (canto inferior direito)", 20, WHITE)
add_para(tf, "5. O browser abre com a cena VR!", 20, WHITE)
add_para(tf, "", 10, WHITE)
add_para(tf, "Se nao tiver 'Go Live': Extensions (Ctrl+Shift+X) > 'Live Server' > Install", 16, ACCENT3)
add_para(tf, "", 6, WHITE)
add_para(tf, "Extensoes uteis: Path Intellisense (Christian Kohler), Color Highlight (Sergii N)", 15, MID_GRAY)


# ── Passo 1-5 (slides condensados) ──
slide = content_slide("Passo 1: Sky (Panorama 360)")
tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
    "Em <a-assets>, carregar a imagem equiretangular:", 18, ACCENT, True)
add_para(tf, '<img id="pano" src="assets/images/panorama.jpg" crossorigin="anonymous" />', 15, ACCENT2, space_before=Pt(8))
add_para(tf, "", 8, WHITE)
add_para(tf, "Na cena:", 18, ACCENT, True)
add_para(tf, '<a-sky src="#pano" rotation="0 -90 0"></a-sky>', 15, ACCENT2, space_before=Pt(8))
add_para(tf, "", 10, WHITE)
add_para(tf, "Resultado: estao DENTRO de uma foto 360!", 18, ACCENT2, True)
add_para(tf, "", 8, WHITE)
add_para(tf, "Termos tecnicos utilizados:", 16, WHITE, True)
add_para(tf, "\u2022 <a-assets> = asset management system (pre-loading)", 15, LIGHT_GRAY)
add_para(tf, '\u2022 id="pano" = DOM identifier (referencia unica)', 15, LIGHT_GRAY)
add_para(tf, "\u2022 <a-sky> = primitivo esfera com radius=5000, material invertido (renderiza por dentro)", 15, LIGHT_GRAY)
add_para(tf, '\u2022 crossorigin="anonymous" = CORS policy (Cross-Origin Resource Sharing)', 15, LIGHT_GRAY)
add_para(tf, "\u2022 src=\"#pano\" = selector CSS que referencia o asset pelo id", 15, LIGHT_GRAY)


slide = content_slide("Passo 2: Formas 3D + EXERCICIO")
tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5),
    "Adicionar formas (primitivos A-Frame):", 18, ACCENT, True)
add_para(tf, "", 4, WHITE)
add_para(tf, '<a-box position="-4 0.7 -5"', 14, ACCENT2)
add_para(tf, '  color="#4CC3D9" class="clickable"></a-box>', 14, ACCENT2)
add_para(tf, "", 6, WHITE)
add_para(tf, '<a-sphere position="3.5 1.5 -5"', 14, ACCENT2)
add_para(tf, '  color="#EF2D5E" radius="1.25"></a-sphere>', 14, ACCENT2)
add_para(tf, "", 6, WHITE)
add_para(tf, '<a-cylinder position="-2 1 -7"', 14, ACCENT2)
add_para(tf, '  color="#FFC65D" height="2"></a-cylinder>', 14, ACCENT2)
add_para(tf, "", 8, WHITE)
add_para(tf, "Atributos-chave:", 16, WHITE, True)
add_para(tf, '\u2022 position="X Y Z" — posicao em metros', 14, LIGHT_GRAY)
add_para(tf, '\u2022 color="#HEX" ou color="nome" — cor', 14, LIGHT_GRAY)
add_para(tf, '\u2022 class="clickable" — torna interativo', 14, LIGHT_GRAY)

add_card(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(5), BG_CARD, ACCENT2)
tf = add_textbox(slide, Inches(7.3), Inches(1.8), Inches(5), Inches(0.5), "EXERCICIO 1", 22, ACCENT2, True)
add_para(tf, "", 6, WHITE)
add_para(tf, "1. Mudem uma cor:", 16, LIGHT_GRAY)
add_para(tf, '   color="#FF0000"  (vermelho)', 14, ACCENT2)
add_para(tf, '   color="#00FF00"  (verde)', 14, ACCENT2)
add_para(tf, '   color="yellow"  (nome tambem funciona)', 14, ACCENT2)
add_para(tf, "", 6, WHITE)
add_para(tf, "2. Mudem uma posicao:", 16, LIGHT_GRAY)
add_para(tf, '   position="0 3 -5"', 14, ACCENT2)
add_para(tf, "", 6, WHITE)
add_para(tf, "3. Copiem um bloco inteiro,", 16, LIGHT_GRAY)
add_para(tf, "   colem abaixo, mudem cor e posicao.", 16, LIGHT_GRAY)
add_para(tf, "", 8, WHITE)
add_para(tf, "Ctrl+S para ver!", 16, ACCENT2, True)


slide = content_slide("Passo 3: Modelo 3D (Coracao)")
tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
    "Em <a-assets>:", 18, ACCENT, True)
add_para(tf, '<a-asset-item id="heart-obj" src="assets/heart/HumanHeart_OBJ.obj"></a-asset-item>', 13, ACCENT2, space_before=Pt(6))
add_para(tf, '<a-asset-item id="heart-mtl" src="assets/heart/HumanHeart_OBJ.mtl"></a-asset-item>', 13, ACCENT2)
add_para(tf, "", 8, WHITE)
add_para(tf, "Na cena:", 18, ACCENT, True)
add_para(tf, '<a-entity obj-model="obj: #heart-obj; mtl: #heart-mtl"', 13, ACCENT2, space_before=Pt(6))
add_para(tf, '  material="side: double" position="0 1.5 -3" scale="0.12 0.12 0.12">', 13, ACCENT2)
add_para(tf, '</a-entity>', 13, ACCENT2)
add_para(tf, "", 8, WHITE)
add_para(tf, "Termos tecnicos:", 16, WHITE, True)
add_para(tf, "\u2022 OBJ = Wavefront Object (formato 3D universal, geometria como vertices e faces)", 14, LIGHT_GRAY)
add_para(tf, "\u2022 MTL = Material Template Library (cores e texturas associadas ao OBJ)", 14, LIGHT_GRAY)
add_para(tf, "\u2022 <a-entity> = entity generica (nao e primitivo — configuramos com componentes)", 14, LIGHT_GRAY)
add_para(tf, '\u2022 obj-model = componente A-Frame que carrega formato OBJ+MTL', 14, LIGHT_GRAY)
add_para(tf, '\u2022 material="side: double" = renderiza AMBAS as faces dos poligonos (backface culling off)', 14, LIGHT_GRAY)
add_para(tf, '\u2022 scale="0.12 0.12 0.12" = 12% do tamanho original (modelos 3D vem em escalas arbitrarias)', 14, LIGHT_GRAY)


slide = content_slide("Passo 4: Animacoes + EXERCICIO")
tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.5),
    "Animacao = atributo HTML declarativo", 18, ACCENT, True)
add_para(tf, "", 6, WHITE)
add_para(tf, "Rotacao continua (cubo):", 15, WHITE, True)
add_para(tf, 'animation="property: rotation;', 13, ACCENT2)
add_para(tf, '  to: 0 390 0; loop: true;', 13, ACCENT2)
add_para(tf, '  dur: 8000; easing: linear"', 13, ACCENT2)
add_para(tf, "", 6, WHITE)
add_para(tf, "Pulso/heartbeat (coracao):", 15, WHITE, True)
add_para(tf, 'animation__pulse="property: scale;', 13, ACCENT3)
add_para(tf, '  from: 0.12 0.12 0.12; to: 0.13 0.13 0.13;', 13, ACCENT3)
add_para(tf, '  dir: alternate; dur: 800; loop: true"', 13, ACCENT3)
add_para(tf, "", 8, WHITE)
add_para(tf, "Termos tecnicos:", 15, WHITE, True)
add_para(tf, "\u2022 property — propriedade CSS/3D a animar", 13, LIGHT_GRAY)
add_para(tf, "\u2022 dur — duration em milissegundos", 13, LIGHT_GRAY)
add_para(tf, "\u2022 easing — curva de aceleracao (linear, easeInOut...)", 13, LIGHT_GRAY)
add_para(tf, "\u2022 dir: alternate — vai e volta (yo-yo)", 13, LIGHT_GRAY)
add_para(tf, "\u2022 animation__nome — multiplas animacoes (__ = sufixo)", 13, LIGHT_GRAY)
add_para(tf, "\u2022 startEvents — dispara com evento (ex: click)", 13, LIGHT_GRAY)

add_card(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(5.2), BG_CARD, ACCENT4)
tf = add_textbox(slide, Inches(7.3), Inches(1.8), Inches(5), Inches(0.5), "EXERCICIO 2", 22, ACCENT4, True)
add_para(tf, "", 6, WHITE)
add_para(tf, "1. Mudem velocidade:", 15, LIGHT_GRAY)
add_para(tf, "   dur: 2000 (rapido)", 14, ACCENT2)
add_para(tf, "   dur: 20000 (lento)", 14, ACCENT2)
add_para(tf, "", 6, WHITE)
add_para(tf, "2. Mudem direcao:", 15, LIGHT_GRAY)
add_para(tf, "   to: 0 -360 0 (anti-horario)", 14, ACCENT2)
add_para(tf, "", 6, WHITE)
add_para(tf, "3. Animem a esfera:", 15, LIGHT_GRAY)
add_para(tf, "   Copiem o animation__pulse", 14, LIGHT_GRAY)
add_para(tf, "", 6, WHITE)
add_para(tf, "4. Mudem easing:", 15, LIGHT_GRAY)
add_para(tf, "   easeInOutBounce", 14, ACCENT2)
add_para(tf, "   easeOutElastic", 14, ACCENT2)
add_para(tf, "", 6, WHITE)
add_para(tf, "Ctrl+S para ver!", 15, ACCENT4, True)


slide = content_slide("Passo 5: Interacoes + Inspector")
tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(5),
    "event-set: interacao SEM JavaScript", 18, ACCENT, True)
add_para(tf, "", 4, WHITE)
add_para(tf, "Hover muda cor:", 15, WHITE, True)
add_para(tf, 'event-set__enter="_event: mouseenter; material.color: #00FFFF"', 13, ACCENT2)
add_para(tf, 'event-set__leave="_event: mouseleave; material.color: #4CC3D9"', 13, ACCENT2)
add_para(tf, "", 6, WHITE)
add_para(tf, "Click muda OUTRO objeto:", 15, WHITE, True)
add_para(tf, 'event-set__click="_event: click; _target: #esfera; material.color: #9B59B6"', 13, ACCENT2)
add_para(tf, "", 6, WHITE)
add_para(tf, "Label ao hover:", 15, WHITE, True)
add_para(tf, 'event-set__enter="_event: mouseenter; _target: #label; visible: true"', 13, ACCENT3)
add_para(tf, "", 8, WHITE)
add_para(tf, "Termos:", 15, WHITE, True)
add_para(tf, "\u2022 _event: evento DOM que dispara a acao", 13, LIGHT_GRAY)
add_para(tf, "\u2022 _target: #id do elemento alvo (se diferente)", 13, LIGHT_GRAY)
add_para(tf, "\u2022 mouseenter/mouseleave = hover in/out", 13, LIGHT_GRAY)
add_para(tf, "\u2022 class='clickable' + raycaster = filtro de interacao", 13, LIGHT_GRAY)

add_card(slide, Inches(7.5), Inches(1.6), Inches(5.2), Inches(5), BG_CARD, ACCENT2)
tf = add_textbox(slide, Inches(7.8), Inches(1.8), Inches(4.6), Inches(0.5), "EXERCICIO 3 FINAL", 22, ACCENT2, True)
add_para(tf, "", 6, WHITE)
add_para(tf, "Adicionem um objeto novo com:", 16, WHITE, True)
add_para(tf, "", 4, WHITE)
add_para(tf, "1. Uma FORMA e COR", 15, LIGHT_GRAY)
add_para(tf, "2. Uma ANIMACAO", 15, LIGHT_GRAY)
add_para(tf, "3. Um efeito HOVER", 15, LIGHT_GRAY)
add_para(tf, "", 8, WHITE)
add_para(tf, "Dica: copiem de objetos existentes!", 15, ACCENT2)
add_para(tf, "", 6, WHITE)
add_para(tf, "Usem o Inspector (Ctrl+Alt+I)", 15, ACCENT, True)
add_para(tf, "para posicionar visualmente!", 15, ACCENT)
add_para(tf, "", 8, WHITE)
add_para(tf, "Tempo: 15 min | Trabalhem em pares!", 16, ACCENT3, True)


# ── Testar VR ──
slide = content_slide("Testar em VR: Meta Quest")
tf = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
    "Como ver no Quest?", 22, ACCENT, True)
add_para(tf, "", 6, WHITE)
add_para(tf, "1. PC e Quest na MESMA rede Wi-Fi", 20, WHITE)
add_para(tf, "2. Live Server mostra: http://192.168.x.x:5500", 20, WHITE)
add_para(tf, "3. No Quest browser, escrever esse IP", 20, WHITE)
add_para(tf, "4. Clicar icone VR (canto inferior direito)", 20, WHITE)
add_para(tf, "5. Thumbstick para andar, olhar para interagir!", 20, WHITE)
add_para(tf, "", 10, WHITE)
add_para(tf, "O cursor (reticulo/crosshair) segue o olhar.", 18, ACCENT2, True)
add_para(tf, "Apontar 1.5s = fuse click (gaze-based interaction).", 18, ACCENT2)


# ══════════════════════════════════════════════════════════════
#  APPS DE REFERENCIA
# ══════════════════════════════════════════════════════════════
slide = content_slide("Referencia: Apps para Criacao de Conteudo VR")
add_table(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2),
    12, 3,
    ["Tipo de Conteudo", "Ferramenta", "Preco / Notas"],
    [
        ["Modelacao 3D", "Blender", "Gratuito (GPL) — recomendado"],
        ["Modelacao 3D", "Autodesk Maya", "~2.016 EUR/ano"],
        ["Scan 3D (mobile)", "Kiri Engine / Polycam", "Gratis (basico) / Pro ~18 EUR/mes"],
        ["Gaussian Splatting", "Luma AI / PostShot / Kiri", "Gratis (basico)"],
        ["Audio espacial", "Audacity + Resonance Audio", "Gratis"],
        ["Imagem 360", "Insta360 app / Google Street View", "Gratis (camera necessaria)"],
        ["Video 360 edicao", "Insta360 Studio / DaVinci Resolve", "Gratis"],
        ["WebVR / WebXR", "A-Frame / Three.js / Babylon.js", "Gratis (open-source)"],
        ["Game Engine", "Unity / Unreal Engine", "Gratis (com limites)"],
        ["AR (sem codigo)", "Adobe Aero / Zappar / Assemblr", "Gratis (basico)"],
        ["VR (sem codigo)", "CoSpaces Edu / Shapes XR", "Gratis (edu) / Pro plans"],
    ])


# ══════════════════════════════════════════════════════════════
#  ENCERRAMENTO
# ══════════════════════════════════════════════════════════════
slide = content_slide("Resumo: O que Aprendemos")
blocks_summary = [
    ("XR e Imersao", "Presenca, cybersickness,\nrequisitos tecnicos,\nevidencia clinica", ACCENT),
    ("Ferramentas", "Som espacial, 360,\n3D, Gaussian Splatting,\nKiri Engine", ACCENT2),
    ("Interatividade", "Game Engines,\nA-Frame + Inspector,\nglossario tecnico", ACCENT3),
    ("Pratica", "Cena VR construida\ndo zero, testada\nno Quest!", ACCENT4),
]
for i, (title, desc, color) in enumerate(blocks_summary):
    x = Inches(0.8 + i * 3.05)
    add_card(slide, x, Inches(1.6), Inches(2.8), Inches(3.8), BG_CARD, color)
    add_textbox(slide, x + Inches(0.3), Inches(1.9), Inches(2.2), Inches(0.5), "\u2713", 28, color, True, PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.3), Inches(2.5), Inches(2.2), Inches(0.5), title, 20, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.3), Inches(3.2), Inches(2.2), Inches(1.8), desc, 14, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

tf = add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1),
    "Proxima aula: usar o scan Kiri da sala como ambiente VR + adicionar mais interacoes!", 18, ACCENT, True)


# ── Recursos com QR codes ──
slide = content_slide("Recursos e Referencias")
add_card(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2), BG_CARD, ACCENT)
tf = add_textbox(slide, Inches(1.1), Inches(1.8), Inches(5), Inches(0.5), "Ferramentas (gratuitas)", 18, ACCENT, True)
items = [
    "A-Frame: aframe.io",
    "A-Frame School: aframe.io/aframe-school",
    "A-Frame Docs: aframe.io/docs",
    "Blender: blender.org",
    "Kiri Engine: kiriengine.app",
    "VS Code: code.visualstudio.com",
    "Three.js: threejs.org",
    "Glitch.com: editor A-Frame online",
]
for item in items:
    add_para(tf, "\u2022 " + item, 13, LIGHT_GRAY, space_before=Pt(5))

add_card(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(5.2), BG_CARD, ACCENT2)
tf = add_textbox(slide, Inches(7.1), Inches(1.8), Inches(5), Inches(0.5), "Assets gratuitos", 18, ACCENT2, True)
items = [
    "Sketchfab: modelos 3D (sketchfab.com)",
    "Poly Haven: HDRIs + texturas (polyhaven.com)",
    "Mixamo: animacoes humanoides (mixamo.com)",
    "Freesound: sons/audio (freesound.org)",
    "Pexels/Unsplash: fotos (pexels.com)",
    "Kenney: assets low-poly (kenney.nl)",
    "Quaternius: modelos gratuitos",
    "ambientCG: materiais PBR (ambientcg.com)",
]
for item in items:
    add_para(tf, "\u2022 " + item, 13, LIGHT_GRAY, space_before=Pt(5))

# QR codes
for url, x_pos, label in [
    ("https://aframe.io", Inches(1.5), "A-Frame"),
    ("https://aframe.io/docs", Inches(4), "Docs"),
    ("https://kiriengine.app", Inches(8), "Kiri Engine"),
    ("https://blender.org", Inches(10.5), "Blender"),
]:
    qr_path = make_qr(url)
    slide.shapes.add_picture(qr_path, x_pos, Inches(5.7), Inches(1.0), Inches(1.0))
    add_textbox(slide, x_pos - Inches(0.2), Inches(6.8), Inches(1.4), Inches(0.3), label, 10, MID_GRAY, alignment=PP_ALIGN.CENTER)


# ── Slide final ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), H)
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
add_textbox(slide, Inches(1.2), Inches(2), Inches(10), Inches(1.5),
    "Obrigado!", 56, WHITE, True, font_name="Segoe UI Semibold")
add_textbox(slide, Inches(1.2), Inches(3.7), Inches(10), Inches(0.8),
    "Questoes? Feedback?", 28, ACCENT, font_name="Segoe UI Light")
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(5), Inches(3), Pt(3))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
add_textbox(slide, Inches(1.2), Inches(5.3), Inches(10), Inches(0.5),
    "Realidade Estendida e Metaverso em Saude — Ferramentas de Criacao de Conteudo", 14, MID_GRAY, font_name="Segoe UI Light")


# ══════════════════════════════════════════════════════════════
#  GUARDAR
# ══════════════════════════════════════════════════════════════
output_path = r"d:\Aula_VR\TesteAulas\Aula_001\apresentacao.pptx"
prs.save(output_path)
print(f"Apresentacao guardada em: {output_path}")
print(f"Total de slides: {len(prs.slides)}")

# Cleanup temp QR files
for f in qr_temp_files:
    try: os.unlink(f)
    except: pass
print("QR codes temp files limpos.")
